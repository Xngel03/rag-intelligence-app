import os
from docx import Document
from pptx import Presentation
import pdfplumber
import pandas as pd
from bs4 import BeautifulSoup

def extract_text(file_path, ext):
    if ext == ".pdf":
        return extract_pdf(file_path)
    elif ext == ".docx":
        return extract_docx(file_path)
    elif ext == ".pptx":
        return extract_pptx(file_path)
    elif ext == ".txt":
        return extract_txt(file_path)
   
    elif ext in [".html", ".htm"]:
        return extract_html(file_path)
    else:
        return "Unsupported file format."

def extract_pdf(file_path):
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
            text += "\n"

            tables = page.extract_tables()
            for table in tables:
                for row in table:
                    text += " | ".join(cell or "" for cell in row) + "\n"
            text += "\n"
    return text

def extract_docx(file_path):
    doc = Document(file_path)
    text = []

    text += [para.text for para in doc.paragraphs]

    for table in doc.tables:
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            text.append(" | ".join(row_data))

    return "\n".join(text)

def extract_pptx(file_path):
    ppt = Presentation(file_path)
    text = ""
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
            elif shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    text += " | ".join(row_data) + "\n"
    return text

def extract_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()



def extract_html(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f, "html.parser")
        return soup.get_text(separator="\n")
